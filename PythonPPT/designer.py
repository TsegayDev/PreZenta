import json
import os
import logging
from typing import Dict, List, Optional, Tuple, Any, Union
from dataclasses import dataclass, field, asdict
from enum import Enum
from pathlib import Path
import colorsys
import hashlib
from datetime import datetime
import warnings
from contextlib import contextmanager

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from lxml import etree

# ============================================================================
# LOGGING CONFIGURATION
# ============================================================================

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ============================================================================
# EXCEPTIONS
# ============================================================================

class PresentationError(Exception):
    """Base exception for presentation generation errors"""
    pass

class LayoutError(PresentationError):
    """Exception raised for layout-related errors"""
    pass

class ThemeError(PresentationError):
    """Exception raised for theme-related errors"""
    pass

class ResourceError(PresentationError):
    """Exception raised for missing or invalid resources"""
    pass

# ============================================================================
# DESIGN SYSTEM & CONFIGURATION
# ============================================================================

class LayoutType(Enum):
    TITLE = "title"
    CONTENT = "content"
    TWO_COLUMN = "two-column"
    THREE_COLUMN = "three-column"
    FOUR_COLUMN = "four-column"
    TWO_COLUMN_HEADED = "two-column-headed"
    THREE_COLUMN_HEADED = "three-column-headed"
    IMAGE_RIGHT = "image-right"
    IMAGE_LEFT = "image-left"
    IMAGE_BACKGROUND = "image-background"
    QUOTE = "quote"
    DATA_VIS = "data-visualization"
    TIMELINE = "timeline"
    COMPARISON = "comparison"
    BULLET_FOCUS = "bullet-focus"
    METRICS = "metrics"
    GRID = "grid"
    HERO = "hero"
    TEAM = "team"
    PROCESS = "process"
    STATS = "stats"

class AnimationType(Enum):
    FADE = "fade"
    SLIDE = "slide"
    ZOOM = "zoom"
    FLIP = "flip"
    NONE = "none"

@dataclass
class Theme:
    """Modern color scheme with dark mode support and dynamic color generation"""
    name: str = "Modern Light"
    background: str = "#F8FAFC"
    surface: str = "#FFFFFF"
    surface_variant: str = "#F1F5F9"
    primary: str = "#3B82F6"
    primary_light: str = "#60A5FA"
    primary_dark: str = "#2563EB"
    secondary: str = "#8B5CF6"
    secondary_light: str = "#A78BFA"
    secondary_dark: str = "#7C3AED"
    accent: str = "#F59E0B"
    success: str = "#10B981"
    warning: str = "#F59E0B"
    error: str = "#EF4444"
    info: str = "#3B82F6"
    text_primary: str = "#0F172A"
    text_secondary: str = "#475569"
    text_muted: str = "#94A3B8"
    border: str = "#E2E8F0"
    shadow: str = "000000"
    
    # Gradient configurations
    primary_gradient: Tuple[str, str] = ("#3B82F6", "#8B5CF6")
    secondary_gradient: Tuple[str, str] = ("#8B5CF6", "#EC4899")
    
    def __post_init__(self):
        """Validate and normalize color values"""
        self._validate_colors()
    
    def _validate_colors(self):
        """Ensure all color values are valid hex codes"""
        for field_name, value in asdict(self).items():
            if isinstance(value, str) and value.startswith('#'):
                if not self._is_valid_hex(value):
                    logger.warning(f"Invalid hex color {value} for {field_name}, using default")
                    setattr(self, field_name, "#000000")
    
    @staticmethod
    def _is_valid_hex(color: str) -> bool:
        """Check if string is a valid hex color"""
        if not color.startswith('#'):
            return False
        color = color.lstrip('#')
        return len(color) in (3, 6) and all(c in '0123456789ABCDEFabcdef' for c in color)
    
    @classmethod
    def dark_mode(cls) -> "Theme":
        return cls(
            name="Modern Dark",
            background="#0F172A",
            surface="#1E293B",
            surface_variant="#334155",
            primary="#60A5FA",
            primary_light="#93C5FD",
            primary_dark="#3B82F6",
            secondary="#A78BFA",
            secondary_light="#C4B5FD",
            secondary_dark="#8B5CF6",
            accent="#FBBF24",
            success="#34D399",
            warning="#FBBF24",
            error="#F87171",
            info="#60A5FA",
            text_primary="#F8FAFC",
            text_secondary="#CBD5E1",
            text_muted="#64748B",
            border="#334155",
            shadow="000000",
            primary_gradient=("#60A5FA", "#A78BFA"),
            secondary_gradient=("#A78BFA", "#F472B6")
        )
    
    @classmethod
    def corporate(cls) -> "Theme":
        return cls(
            name="Corporate",
            background="#FFFFFF",
            surface="#F8FAFC",
            surface_variant="#F1F5F9",
            primary="#1E40AF",
            primary_light="#3B82F6",
            primary_dark="#1E3A8A",
            secondary="#64748B",
            secondary_light="#94A3B8",
            secondary_dark="#475569",
            accent="#F59E0B",
            success="#059669",
            warning="#D97706",
            error="#DC2626",
            info="#2563EB",
            text_primary="#111827",
            text_secondary="#4B5563",
            text_muted="#9CA3AF",
            border="#E5E7EB"
        )
    
    @classmethod
    def minimalist(cls) -> "Theme":
        return cls(
            name="Minimalist",
            background="#FFFFFF",
            surface="#FAFAFA",
            surface_variant="#F5F5F5",
            primary="#000000",
            primary_light="#333333",
            primary_dark="#000000",
            secondary="#666666",
            secondary_light="#999999",
            secondary_dark="#333333",
            accent="#757575",
            success="#4CAF50",
            warning="#FF9800",
            error="#F44336",
            info="#2196F3",
            text_primary="#212121",
            text_secondary="#757575",
            text_muted="#BDBDBD",
            border="#EEEEEE"
        )

@dataclass
class Typography:
    """Typography scale following modern design systems with responsive sizing"""
    hero: int = 54
    h1: int = 40
    h2: int = 32
    h3: int = 24
    h4: int = 20
    body_large: int = 18
    body: int = 14
    caption: int = 12
    tiny: int = 10
    font_family: str = "Calibri"
    font_family_heading: str = "Calibri Light"
    font_family_mono: str = "Consolas"
    line_height_body: float = 1.5
    line_height_heading: float = 1.2
    
    def get_size(self, level: str) -> int:
        """Get font size by level name"""
        return getattr(self, level, self.body)

@dataclass
class Spacing:
    """Consistent spacing system with responsive scaling"""
    unit: float = 0.125  # Base unit in inches (approx 9px)
    
    @property
    def xs(self) -> float:
        return self.unit * 1
    
    @property
    def sm(self) -> float:
        return self.unit * 2
    
    @property
    def md(self) -> float:
        return self.unit * 4
    
    @property
    def lg(self) -> float:
        return self.unit * 6
    
    @property
    def xl(self) -> float:
        return self.unit * 8
    
    @property
    def xxl(self) -> float:
        return self.unit * 12
    
    def scale(self, factor: float) -> 'Spacing':
        """Create scaled version of spacing"""
        new_spacing = Spacing()
        new_spacing.unit = self.unit * factor
        return new_spacing

# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor with validation and error handling"""
    try:
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 6:
            return RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
        elif len(hex_color) == 3:
            return RGBColor(
                int(hex_color[0]*2, 16),
                int(hex_color[1]*2, 16),
                int(hex_color[2]*2, 16)
            )
    except (ValueError, TypeError) as e:
        logger.error(f"Error converting hex color {hex_color}: {e}")
    return RGBColor(0, 0, 0)

def rgb_to_hex(rgb: RGBColor) -> str:
    """Convert RGBColor to hex string"""
    return f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"

def adjust_brightness(hex_color: str, factor: float) -> str:
    """Adjust brightness of hex color (factor: -1 to 1) with bounds checking"""
    try:
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16) / 255.0
        g = int(hex_color[2:4], 16) / 255.0
        b = int(hex_color[4:6], 16) / 255.0
        
        h, s, l = colorsys.rgb_to_hls(r, g, b)
        l = max(0, min(1, l + factor))
        r, g, b = colorsys.hls_to_rgb(h, l, s)
        
        return f"#{int(r*255):02x}{int(g*255):02x}{int(b*255):02x}"
    except Exception as e:
        logger.error(f"Error adjusting brightness: {e}")
        return hex_color

def blend_colors(color1: str, color2: str, ratio: float = 0.5) -> str:
    """Blend two colors with given ratio"""
    try:
        rgb1 = hex_to_rgb(color1)
        rgb2 = hex_to_rgb(color2)
        
        r = int(rgb1.r * (1 - ratio) + rgb2.r * ratio)
        g = int(rgb1.g * (1 - ratio) + rgb2.g * ratio)
        b = int(rgb1.b * (1 - ratio) + rgb2.b * ratio)
        
        return f"#{r:02x}{g:02x}{b:02x}"
    except Exception:
        return color1

def get_contrast_color(hex_color: str) -> str:
    """Get contrasting text color (black or white) for background"""
    try:
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16) / 255.0
        g = int(hex_color[2:4], 16) / 255.0
        b = int(hex_color[4:6], 16) / 255.0
        
        # Calculate luminance
        luminance = 0.299 * r + 0.587 * g + 0.114 * b
        
        return "#000000" if luminance > 0.5 else "#FFFFFF"
    except Exception:
        return "#000000"

@contextmanager
def safe_xml_operation():
    """Context manager for safe XML operations"""
    try:
        yield
    except Exception as e:
        logger.error(f"XML operation failed: {e}")

def set_shape_rounded_corners(shape, radius: float = 0.1):
    """Apply rounded corners to a shape using XML manipulation with error handling"""
    with safe_xml_operation():
        try:
            sp = shape._element
            spPr = sp.find(qn('a:spPr'))
            if spPr is not None:
                # Remove existing geometry
                for geom in spPr.findall(qn('a:custGeom')):
                    spPr.remove(geom)
                for geom in spPr.findall(qn('a:prstGeom')):
                    spPr.remove(geom)
                
                # Apply rounded rectangle
                from pptx.oxml.shapes import CT_PresetGeometry2D
                prstGeom = CT_PresetGeometry2D.new('roundRect')
                
                # Set corner radius
                avLst = etree.Element(qn('a:avLst'))
                gd = etree.Element(qn('a:gd'))
                gd.set('name', 'adj')
                gd.set('fmla', f'val {int(radius * 100000)}')
                avLst.append(gd)
                prstGeom.append(avLst)
                
                spPr.append(prstGeom)
        except Exception as e:
            logger.debug(f"Failed to set rounded corners: {e}")

# ============================================================================
# CACHE MANAGER
# ============================================================================

class ResourceCache:
    """Cache for frequently used resources"""
    
    def __init__(self, max_size: int = 100):
        self.cache = {}
        self.max_size = max_size
    
    def get(self, key: str) -> Any:
        """Get item from cache"""
        return self.cache.get(key)
    
    def set(self, key: str, value: Any):
        """Set item in cache with LRU-like behavior"""
        if len(self.cache) >= self.max_size:
            # Remove oldest item
            oldest = next(iter(self.cache))
            del self.cache[oldest]
        self.cache[key] = value
    
    def get_key(self, *args, **kwargs) -> str:
        """Generate cache key from arguments"""
        key_string = str(args) + str(sorted(kwargs.items()))
        return hashlib.md5(key_string.encode()).hexdigest()

# ============================================================================
# VALIDATORS
# ============================================================================

class SlideDataValidator:
    """Validate slide data structure"""
    
    @staticmethod
    def validate_slide(slide_data: Dict) -> List[str]:
        """Validate slide data and return list of errors"""
        errors = []
        
        if not isinstance(slide_data, dict):
            errors.append("Slide data must be a dictionary")
            return errors
        
        # Check required fields
        if 'layout' not in slide_data:
            errors.append("Missing required field: 'layout'")
        
        # Validate layout type
        layout = slide_data.get('layout')
        if layout and not any(layout == lt.value for lt in LayoutType):
            errors.append(f"Invalid layout type: {layout}")
        
        # Validate based on layout
        if layout == 'title':
            if 'title' not in slide_data:
                errors.append("Title slide missing 'title' field")
        
        elif layout in ['two-column', 'three-column', 'four-column']:
            if 'columns' not in slide_data:
                errors.append(f"{layout} slide missing 'columns' field")
            elif not isinstance(slide_data['columns'], list):
                errors.append("'columns' must be a list")
        
        return errors

# ============================================================================
# MAIN GENERATOR CLASS
# ============================================================================

class ModernPresentationGenerator:
    """
    Advanced PowerPoint generator with modern design capabilities:
    - Glassmorphism & neumorphism effects
    - Data visualization layouts
    - Smooth animations and transitions
    - Responsive grid systems
    - Dark/Light theme support
    - Resource caching
    - Error recovery
    """
    
    def __init__(self, json_data: Dict[str, Any], output_dir: str = "."):
        self.data = json_data
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Initialize presentation
        self.prs = Presentation()
        
        # Setup slide dimensions (16:9 standard)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # Load configuration with error handling
        try:
            theme_config = self.data.get('theme', {})
            self.theme = self._load_theme(theme_config)
            self.typography = self._load_typography(theme_config.get('typography', {}))
            self.spacing = Spacing()
            
            # Scale spacing if specified
            spacing_scale = theme_config.get('spacing_scale', 1.0)
            self.spacing = self.spacing.scale(spacing_scale)
            
        except Exception as e:
            logger.error(f"Error loading configuration: {e}")
            raise ThemeError(f"Failed to load theme: {e}")
        
        # Track slide numbers
        self.current_slide_num = 0
        self.total_slides = len(self.data.get('slides', []))
        
        # Initialize cache
        self.cache = ResourceCache()
        
        # Initialize validator
        self.validator = SlideDataValidator()
        
        # Track errors for reporting
        self.errors = []
        self.warnings = []
    
    def _load_theme(self, config: Dict) -> Theme:
        """Load theme from config or use default with error recovery"""
        theme_name = config.get('name', 'light')
        
        try:
            if theme_name == 'dark':
                base = Theme.dark_mode()
            elif theme_name == 'corporate':
                base = Theme.corporate()
            elif theme_name == 'minimalist':
                base = Theme.minimalist()
            else:
                base = Theme()
            
            # Override with any custom colors provided
            for key, value in config.items():
                if hasattr(base, key):
                    if key.endswith('_gradient') and isinstance(value, list):
                        setattr(base, key, tuple(value[:2]))
                    elif isinstance(value, str):
                        setattr(base, key, value)
            
            return base
            
        except Exception as e:
            logger.error(f"Error loading theme, using default: {e}")
            self.warnings.append(f"Theme loading failed, using default: {e}")
            return Theme()
    
    def _load_typography(self, config: Dict) -> Typography:
        """Load typography settings with defaults"""
        try:
            return Typography(
                hero=config.get('hero', 54),
                h1=config.get('h1', 40),
                h2=config.get('h2', 32),
                h3=config.get('h3', 24),
                h4=config.get('h4', 20),
                body_large=config.get('body_large', 18),
                body=config.get('body', 14),
                caption=config.get('caption', 12),
                tiny=config.get('tiny', 10),
                font_family=config.get('font_family', 'Calibri'),
                font_family_heading=config.get('font_family_heading', 'Calibri Light'),
                font_family_mono=config.get('font_family_mono', 'Consolas'),
                line_height_body=config.get('line_height_body', 1.5),
                line_height_heading=config.get('line_height_heading', 1.2)
            )
        except Exception as e:
            logger.error(f"Error loading typography: {e}")
            return Typography()
    
    def generate(self, output_filename: str = None) -> str:
        """Generate the complete presentation with error recovery"""
        if output_filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"Presentation_{timestamp}.pptx"
        
        output_path = self.output_dir / output_filename
        
        logger.info(f"Generating {self.total_slides} slides...")
        
        # Validate all slides first
        for i, slide_data in enumerate(self.data.get('slides', [])):
            errors = self.validator.validate_slide(slide_data)
            for error in errors:
                self.errors.append(f"Slide {i+1}: {error}")
        
        if self.errors:
            logger.error(f"Found {len(self.errors)} validation errors")
            for error in self.errors[:5]:  # Show first 5 errors
                logger.error(f"  {error}")
            if len(self.errors) > 5:
                logger.error(f"  ... and {len(self.errors) - 5} more")
            
            if self.data.get('settings', {}).get('strict_mode', False):
                raise PresentationError("Validation failed in strict mode")
        
        # Generate slides with error recovery
        for i, slide_data in enumerate(self.data.get('slides', [])):
            self.current_slide_num = i + 1
            try:
                logger.info(f"Creating slide {self.current_slide_num}/{self.total_slides}: {slide_data.get('title', 'Untitled')[:40]}...")
                self._create_slide(slide_data)
            except Exception as e:
                logger.error(f"Failed to create slide {self.current_slide_num}: {e}")
                self.warnings.append(f"Slide {self.current_slide_num} failed: {e}")
                
                # Create error slide as fallback
                self._create_error_slide(slide_data, str(e))
        
        # Add slide transitions if specified
        if self.data.get('settings', {}).get('transitions', False):
            self._add_transitions()
        
        # Save presentation
        try:
            self.prs.save(str(output_path))
            logger.info(f"✓ Presentation saved successfully to: {output_path}")
            logger.info(f"  Theme: {self.theme.name}")
            logger.info(f"  Dimensions: 16:9 (13.33\" x 7.5\")")
            
            if self.warnings:
                logger.warning(f"  Warnings: {len(self.warnings)}")
            
            return str(output_path)
            
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise PresentationError(f"Save failed: {e}")
    
    def _create_error_slide(self, original_data: Dict, error_message: str):
        """Create a fallback slide when normal rendering fails"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        self._set_slide_background(slide)
        
        # Error message
        self._add_text_box(
            slide, "⚠️ Slide Generation Error",
            0.6, 2.0, 12.1, 0.8,
            font_size=self.typography.h2,
            font_color=self.theme.error,
            bold=True,
            align=PP_ALIGN.CENTER
        )
        
        self._add_text_box(
            slide, f"Error: {error_message}",
            0.6, 3.0, 12.1, 1.0,
            font_size=self.typography.body,
            font_color=self.theme.text_secondary,
            align=PP_ALIGN.CENTER
        )
        
        self._add_text_box(
            slide, f"Original title: {original_data.get('title', 'Unknown')}",
            0.6, 4.0, 12.1, 0.5,
            font_size=self.typography.caption,
            font_color=self.theme.text_muted,
            align=PP_ALIGN.CENTER
        )
    
    def _create_slide(self, slide_data: Dict):
        """Route to appropriate layout renderer"""
        layout_type = slide_data.get('layout', 'content')
        
        # Use blank layout (index 6)
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background
        self._set_slide_background(slide)
        
        # Route to specific renderer
        renderers = {
            'title': self._render_title_slide,
            'content': self._render_content_slide,
            'two-column': lambda s, d: self._render_columns(s, d, 2),
            'three-column': lambda s, d: self._render_columns(s, d, 3),
            'four-column': lambda s, d: self._render_columns(s, d, 4),
            'two-column-headed': lambda s, d: self._render_columns(s, d, 2, headed=True),
            'three-column-headed': lambda s, d: self._render_columns(s, d, 3, headed=True),
            'image-right': lambda s, d: self._render_image_split(s, d, 'right'),
            'image-left': lambda s, d: self._render_image_split(s, d, 'left'),
            'image-background': self._render_image_background,
            'quote': self._render_quote_slide,
            'data-visualization': self._render_data_slide,
            'timeline': self._render_timeline_slide,
            'comparison': self._render_comparison_slide,
            'bullet-focus': self._render_bullet_focus_slide,
            'metrics': self._render_metrics_slide,
            'grid': self._render_grid_slide,
            'hero': self._render_hero_slide,
            'team': self._render_team_slide,
            'process': self._render_process_slide,
            'stats': self._render_stats_slide
        }
        
        renderer = renderers.get(layout_type)
        if renderer:
            try:
                renderer(slide, slide_data)
            except Exception as e:
                logger.error(f"Error in renderer {layout_type}: {e}")
                raise
        else:
            logger.warning(f"Unknown layout type: {layout_type}, using content layout")
            self._render_content_slide(slide, slide_data)
        
        # Add slide number if enabled
        if self.data.get('settings', {}).get('slide_numbers', True):
            self._add_slide_number(slide)
    
    def _set_slide_background(self, slide):
        """Apply theme background with optional gradient"""
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(self.theme.background)
        except Exception as e:
            logger.error(f"Failed to set slide background: {e}")
    
    def _add_text_box(self, slide, text: str, left: float, top: float, 
                      width: float, height: float, font_size: int = None,
                      font_color: str = None, bold: bool = False, 
                      align: PP_ALIGN = PP_ALIGN.LEFT, 
                      font_family: str = None, italic: bool = False,
                      line_spacing: float = 1.2, underline: bool = False,
                      font_file: str = None) -> Optional[Any]:
        """Enhanced text box with typography controls and error handling"""
        try:
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                              Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            # Set line spacing
            tf.paragraphs[0].line_spacing = line_spacing
            
            p = tf.paragraphs[0]
            p.alignment = align
            
            # Clear any existing runs
            p.clear()
            
            # Add text if provided
            if text:
                run = p.add_run()
                run.text = str(text)
                
                # Apply formatting
                run.font.size = Pt(font_size or self.typography.body)
                run.font.color.rgb = hex_to_rgb(font_color or self.theme.text_primary)
                run.font.bold = bold
                run.font.italic = italic
                run.font.underline = underline
                run.font.name = font_family or self.typography.font_family
            
            return txBox
            
        except Exception as e:
            logger.error(f"Failed to add text box: {e}")
            return None
    
    def _add_glass_card(self, slide, left: float, top: float, 
                        width: float, height: float, title: str = None,
                        content: str = None, accent_color: str = None,
                        opacity: float = 0.95, rounded: bool = True,
                        shadow: bool = True) -> Optional[Any]:
        """Create a glassmorphism-style card with enhanced effects"""
        try:
            # Main card shape with rounded corners
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
            card = slide.shapes.add_shape(
                shape_type,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            
            # Style the card
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb(self.theme.surface)
            card.line.color.rgb = hex_to_rgb(self.theme.border)
            card.line.width = Pt(1)
            
            # Add shadow effect if requested
            if shadow:
                self._add_shadow(slide, card, left, top, width, height)
            
            # Add accent line if specified
            if accent_color:
                accent = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left), Inches(top),
                    Inches(0.08), Inches(height)
                )
                accent.fill.solid()
                accent.fill.fore_color.rgb = hex_to_rgb(accent_color)
                accent.line.fill.background()
            
            # Add content
            padding = 0.25
            current_top = top + padding
            
            if title:
                self._add_text_box(
                    slide, title,
                    left + padding + (0.15 if accent_color else 0),
                    current_top,
                    width - (padding * 2) - (0.15 if accent_color else 0),
                    0.6,
                    font_size=self.typography.h3,
                    font_color=self.theme.text_primary,
                    bold=True,
                    font_family=self.typography.font_family_heading
                )
                current_top += 0.7
            
            if content:
                self._add_text_box(
                    slide, content,
                    left + padding + (0.15 if accent_color else 0),
                    current_top,
                    width - (padding * 2) - (0.15 if accent_color else 0),
                    height - current_top + top - padding - 0.1,
                    font_size=self.typography.body,
                    font_color=self.theme.text_secondary,
                    line_spacing=self.typography.line_height_body
                )
            
            return card
            
        except Exception as e:
            logger.error(f"Failed to create glass card: {e}")
            return None
    
    def _add_shadow(self, slide, card, left: float, top: float, 
                    width: float, height: float):
        """Add subtle shadow effect to a shape"""
        try:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + 0.03), Inches(top + 0.03),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
            shadow.fill.transparency = 0.7
            shadow.line.fill.background()
            
            # Move shadow behind card
            spTree = slide.shapes._spTree
            sp = shadow._element
            spTree.remove(sp)
            spTree.insert(2, sp)
            
        except Exception as e:
            logger.debug(f"Failed to add shadow: {e}")
    
    def _add_decorative_line(self, slide, left: float, top: float,
                            width: float, color: str = None, 
                            height: float = 0.06, dashed: bool = False) -> Optional[Any]:
        """Add a decorative accent line"""
        try:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = hex_to_rgb(color or self.theme.primary)
            line.line.fill.background()
            
            if dashed:
                line.line.dash_style = 2  # Dashed line
            
            return line
            
        except Exception as e:
            logger.error(f"Failed to add decorative line: {e}")
            return None
    
    def _render_title_slide(self, slide, data: Dict):
        """Modern title slide with gradient accent and typography hierarchy"""
        title = data.get('title', '')
        subtitle = data.get('subtitle', data.get('content', ''))
        author = data.get('author', '')
        date = data.get('date', '')
        logo = data.get('logo', '')
        
        try:
            # Decorative top bar with gradient effect
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.333), Inches(0.15)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = hex_to_rgb(self.theme.primary)
            top_bar.line.fill.background()
            
            # Main title
            self._add_text_box(
                slide, title,
                0.8, 2.8, 11.7, 1.2,
                font_size=self.typography.hero,
                font_color=self.theme.text_primary,
                bold=True,
                align=PP_ALIGN.CENTER,
                font_family=self.typography.font_family_heading
            )
            
            # Decorative line under title
            line_width = min(len(title) * 0.15, 4)
            line_left = (13.333 - line_width) / 2
            self._add_decorative_line(slide, line_left, 4.1, line_width, 
                                     self.theme.accent, 0.08)
            
            # Subtitle
            if subtitle:
                self._add_text_box(
                    slide, subtitle,
                    1.0, 4.4, 11.3, 0.8,
                    font_size=self.typography.h3,
                    font_color=self.theme.text_secondary,
                    align=PP_ALIGN.CENTER
                )
            
            # Footer info
            footer_text = []
            if author:
                footer_text.append(author)
            if date:
                footer_text.append(date)
            
            if footer_text:
                self._add_text_box(
                    slide, " | ".join(footer_text),
                    0.5, 6.8, 12.3, 0.4,
                    font_size=self.typography.caption,
                    font_color=self.theme.text_muted,
                    align=PP_ALIGN.CENTER
                )
            
            # Bottom accent
            bottom_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(7.35),
                Inches(13.333), Inches(0.15)
            )
            bottom_bar.fill.solid()
            bottom_bar.fill.fore_color.rgb = hex_to_rgb(self.theme.secondary)
            bottom_bar.line.fill.background()
            
            # Add logo if provided
            if logo and os.path.exists(logo):
                try:
                    slide.shapes.add_picture(
                        logo,
                        Inches(6.0), Inches(0.8),
                        Inches(1.333), Inches(1.333)
                    )
                except Exception as e:
                    logger.warning(f"Failed to add logo: {e}")
                    
        except Exception as e:
            logger.error(f"Error in title slide: {e}")
            raise
    
    def _render_content_slide(self, slide, data: Dict):
        """Standard content slide with glass card"""
        title = data.get('title', '')
        content = data.get('content', '')
        bullets = data.get('bullets', [])
        
        try:
            # Header section
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.9,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                font_family=self.typography.font_family_heading
            )
            
            # Decorative line
            self._add_decorative_line(slide, 0.6, 1.3, 1.5, self.theme.accent)
            
            # Prepare content text
            if bullets:
                content_text = "\n".join([f"• {bullet}" for bullet in bullets])
            else:
                content_text = content
            
            # Content card
            self._add_glass_card(
                slide, 0.6, 1.6, 12.1, 5.4,
                content=content_text,
                accent_color=self.theme.primary
            )
            
        except Exception as e:
            logger.error(f"Error in content slide: {e}")
            raise
    
    def _render_columns(self, slide, data: Dict, num_cols: int, headed: bool = False):
        """Multi-column layout with consistent spacing"""
        title = data.get('title', '')
        columns = data.get('columns', [])
        
        try:
            # Header
            if title:
                self._add_text_box(
                    slide, title,
                    0.6, 0.5, 12.1, 0.8,
                    font_size=self.typography.h1,
                    font_color=self.theme.text_primary,
                    bold=True,
                    align=PP_ALIGN.CENTER if num_cols > 1 else PP_ALIGN.LEFT,
                    font_family=self.typography.font_family_heading
                )
            
            # Calculate column dimensions
            margin = 0.6
            gap = 0.4
            total_width = 13.333 - (margin * 2)
            col_width = (total_width - (gap * (num_cols - 1))) / num_cols
            top = 1.5 if title else 0.8
            height = 5.6 if title else 6.3
            
            # Render columns
            for i, col in enumerate(columns[:num_cols]):
                left = margin + (i * (col_width + gap))
                accent = self.theme.accent if headed and i == 0 else None
                
                self._add_glass_card(
                    slide, left, top, col_width, height,
                    title=col.get('title', ''),
                    content=col.get('content', ''),
                    accent_color=accent
                )
                
        except Exception as e:
            logger.error(f"Error in columns slide: {e}")
            raise
    
    def _render_image_split(self, slide, data: Dict, position: str):
        """Two-column layout with image and text"""
        title = data.get('title', '')
        content = data.get('content', '')
        image_path = data.get('image', '')
        
        try:
            text_left = 0.6 if position == 'left' else 7.0
            img_left = 7.0 if position == 'left' else 0.6
            
            # Text section
            self._add_text_box(
                slide, title,
                text_left, 1.2, 5.7, 0.8,
                font_size=self.typography.h2,
                font_color=self.theme.text_primary,
                bold=True,
                font_family=self.typography.font_family_heading
            )
            
            self._add_decorative_line(slide, text_left, 2.0, 1.0, self.theme.primary)
            
            self._add_glass_card(
                slide, text_left, 2.3, 5.7, 4.5,
                content=content
            )
            
            # Image section
            if image_path and os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(img_left), Inches(1.2),
                        Inches(5.7), Inches(5.8)
                    )
                except Exception as e:
                    logger.warning(f"Failed to add image {image_path}: {e}")
                    self._add_image_placeholder(slide, img_left, 1.2, 5.7, 5.8)
            else:
                self._add_image_placeholder(slide, img_left, 1.2, 5.7, 5.8)
                
        except Exception as e:
            logger.error(f"Error in image split slide: {e}")
            raise
    
    def _add_image_placeholder(self, slide, left, top, width, height):
        """Add a styled placeholder for images"""
        try:
            # Background shape
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = hex_to_rgb(self.theme.surface_variant)
            placeholder.line.color.rgb = hex_to_rgb(self.theme.border)
            placeholder.line.width = Pt(2)
            placeholder.line.dash_style = 2  # Dashed line
            
            # Icon/text indicating placeholder
            self._add_text_box(
                slide, "📷 Image Placeholder",
                left, top + height/2 - 0.3, width, 0.6,
                font_size=self.typography.body,
                font_color=self.theme.text_muted,
                align=PP_ALIGN.CENTER
            )
            
        except Exception as e:
            logger.error(f"Failed to add image placeholder: {e}")
    
    def _render_image_background(self, slide, data: Dict):
        """Full-bleed image background with text overlay"""
        title = data.get('title', '')
        content = data.get('content', '')
        image_path = data.get('image', '')
        overlay_opacity = data.get('overlay_opacity', 0.4)
        
        try:
            # Add image as background
            if image_path and os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(0), Inches(0),
                        Inches(13.333), Inches(7.5)
                    )
                except Exception as e:
                    logger.warning(f"Failed to add background image: {e}")
            
            # Dark overlay for text readability
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.333), Inches(7.5)
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = overlay_opacity
            overlay.line.fill.background()
            
            # Send overlay to back
            spTree = slide.shapes._spTree
            sp = overlay._element
            spTree.remove(sp)
            spTree.insert(2, sp)
            
            # Content
            self._add_text_box(
                slide, title,
                1.0, 2.5, 11.3, 1.0,
                font_size=self.typography.hero,
                font_color="#FFFFFF",
                bold=True,
                align=PP_ALIGN.CENTER,
                font_family=self.typography.font_family_heading
            )
            
            self._add_text_box(
                slide, content,
                1.5, 3.8, 10.3, 1.5,
                font_size=self.typography.body_large,
                font_color="#E2E8F0",
                align=PP_ALIGN.CENTER,
                line_spacing=self.typography.line_height_body
            )
            
        except Exception as e:
            logger.error(f"Error in image background slide: {e}")
            raise
    
    def _render_quote_slide(self, slide, data: Dict):
        """Stylized quote layout"""
        quote = data.get('quote', data.get('content', ''))
        author = data.get('author', '')
        role = data.get('role', '')
        
        try:
            # Large quote mark decoration
            self._add_text_box(
                slide, '"',
                0.8, 1.5, 1.0, 1.0,
                font_size=72,
                font_color=self.theme.primary_light,
                bold=True
            )
            
            # Quote text
            self._add_text_box(
                slide, quote,
                1.8, 2.0, 9.7, 3.0,
                font_size=self.typography.h2,
                font_color=self.theme.text_primary,
                italic=True,
                align=PP_ALIGN.LEFT,
                line_spacing=self.typography.line_height_heading * 1.2
            )
            
            # Author line
            if author:
                author_text = f"— {author}"
                if role:
                    author_text += f", {role}"
                
                self._add_text_box(
                    slide, author_text,
                    1.8, 5.2, 9.7, 0.5,
                    font_size=self.typography.body_large,
                    font_color=self.theme.text_secondary,
                    bold=True
                )
            
            # Decorative line
            self._add_decorative_line(slide, 1.8, 5.8, 2.0, self.theme.accent)
            
        except Exception as e:
            logger.error(f"Error in quote slide: {e}")
            raise
    
    def _render_data_slide(self, slide, data: Dict):
        """Data visualization slide with chart placeholder"""
        title = data.get('title', '')
        chart_type = data.get('chart_type', 'bar')
        chart_data = data.get('chart_data', {})
        
        try:
            # Header
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.8,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                font_family=self.typography.font_family_heading
            )
            
            # Chart placeholder area
            chart_left = 0.6
            chart_top = 1.5
            chart_width = 8.0
            chart_height = 5.5
            
            # Chart background
            chart_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(chart_left), Inches(chart_top),
                Inches(chart_width), Inches(chart_height)
            )
            chart_bg.fill.solid()
            chart_bg.fill.fore_color.rgb = hex_to_rgb(self.theme.surface)
            chart_bg.line.color.rgb = hex_to_rgb(self.theme.border)
            
            # Try to generate simple chart visualization
            if chart_data and 'values' in chart_data:
                self._render_simple_chart(slide, chart_left, chart_top, 
                                        chart_width, chart_height, 
                                        chart_type, chart_data)
            else:
                # Chart placeholder text
                self._add_text_box(
                    slide, f"[{chart_type.upper()} CHART]\nInsert chart data here",
                    chart_left + 0.5, chart_top + 2.5, chart_width - 1.0, 1.0,
                    font_size=self.typography.body,
                    font_color=self.theme.text_muted,
                    align=PP_ALIGN.CENTER
                )
            
            # Key metrics side panel
            metrics = data.get('metrics', [])
            if metrics:
                panel_left = 9.0
                self._add_glass_card(
                    slide, panel_left, 1.5, 3.7, 5.5,
                    title="Key Metrics",
                    content="\n\n".join([f"• {m}" for m in metrics]),
                    accent_color=self.theme.primary
                )
                
        except Exception as e:
            logger.error(f"Error in data slide: {e}")
            raise
    
    def _render_simple_chart(self, slide, left, top, width, height, chart_type, data):
        """Render a simple chart visualization"""
        try:
            values = data.get('values', [])
            labels = data.get('labels', [])
            
            if not values:
                return
            
            max_value = max(values)
            chart_left = left + 0.5
            chart_top = top + 0.5
            chart_width = width - 1.0
            chart_height = height - 1.0
            
            if chart_type == 'bar':
                bar_width = chart_width / len(values) * 0.6
                bar_spacing = chart_width / len(values) * 0.4
                
                for i, value in enumerate(values):
                    bar_height = (value / max_value) * chart_height * 0.8
                    bar_x = chart_left + i * (bar_width + bar_spacing)
                    bar_y = chart_top + chart_height * 0.9 - bar_height
                    
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(bar_x), Inches(bar_y),
                        Inches(bar_width), Inches(bar_height)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = hex_to_rgb(self.theme.primary)
                    bar.line.fill.background()
                    
                    # Add value label
                    if labels and i < len(labels):
                        self._add_text_box(
                            slide, labels[i],
                            bar_x - 0.2, bar_y - 0.3, bar_width + 0.4, 0.3,
                            font_size=self.typography.caption,
                            font_color=self.theme.text_muted,
                            align=PP_ALIGN.CENTER
                        )
            
        except Exception as e:
            logger.error(f"Failed to render chart: {e}")
    
    def _render_timeline_slide(self, slide, data: Dict):
        """Horizontal timeline layout with improved visualization"""
        title = data.get('title', '')
        events = data.get('events', [])
        
        try:
            # Header
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.8,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                font_family=self.typography.font_family_heading
            )
            
            if not events:
                return
            
            # Timeline line
            timeline_y = 4.0
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.0), Inches(timeline_y),
                Inches(11.3), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = hex_to_rgb(self.theme.primary)
            line.line.fill.background()
            
            # Event nodes
            num_events = len(events)
            spacing = 11.3 / (num_events - 1) if num_events > 1 else 5.65
            
            for i, event in enumerate(events):
                x = 1.0 + (i * spacing)
                
                # Node circle
                node = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x - 0.15), Inches(timeline_y - 0.125),
                    Inches(0.3), Inches(0.3)
                )
                node.fill.solid()
                node.fill.fore_color.rgb = hex_to_rgb(self.theme.primary)
                node.line.fill.background()
                
                # Date label
                self._add_text_box(
                    slide, event.get('date', ''),
                    x - 1.0, timeline_y - 1.0, 2.0, 0.4,
                    font_size=self.typography.caption,
                    font_color=self.theme.text_muted,
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
                # Event title
                self._add_text_box(
                    slide, event.get('title', ''),
                    x - 1.0, timeline_y + 0.3, 2.0, 0.6,
                    font_size=self.typography.body,
                    font_color=self.theme.text_primary,
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
                # Description
                self._add_text_box(
                    slide, event.get('desc', ''),
                    x - 1.0, timeline_y + 0.9, 2.0, 1.5,
                    font_size=self.typography.caption,
                    font_color=self.theme.text_secondary,
                    align=PP_ALIGN.CENTER,
                    line_spacing=self.typography.line_height_body
                )
                
        except Exception as e:
            logger.error(f"Error in timeline slide: {e}")
            raise
    
    def _render_comparison_slide(self, slide, data: Dict):
        """Side-by-side comparison layout"""
        title = data.get('title', '')
        left_title = data.get('left_title', 'Option A')
        right_title = data.get('right_title', 'Option B')
        left_points = data.get('left_points', [])
        right_points = data.get('right_points', [])
        
        try:
            # Header
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.8,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                align=PP_ALIGN.CENTER,
                font_family=self.typography.font_family_heading
            )
            
            # VS indicator in center
            vs_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(6.166), Inches(3.0),
                Inches(1.0), Inches(1.0)
            )
            vs_circle.fill.solid()
            vs_circle.fill.fore_color.rgb = hex_to_rgb(self.theme.accent)
            vs_circle.line.fill.background()
            
            self._add_text_box(
                slide, "VS",
                6.166, 3.25, 1.0, 0.5,
                font_size=self.typography.h3,
                font_color="#FFFFFF",
                bold=True,
                align=PP_ALIGN.CENTER
            )
            
            # Left column
            left_content = "\n".join([f"✓ {p}" for p in left_points])
            self._add_glass_card(
                slide, 0.6, 1.5, 5.0, 5.5,
                title=left_title,
                content=left_content,
                accent_color=self.theme.primary
            )
            
            # Right column
            right_content = "\n".join([f"✓ {p}" for p in right_points])
            self._add_glass_card(
                slide, 7.7, 1.5, 5.0, 5.5,
                title=right_title,
                content=right_content,
                accent_color=self.theme.secondary
            )
            
        except Exception as e:
            logger.error(f"Error in comparison slide: {e}")
            raise
    
    def _render_bullet_focus_slide(self, slide, data: Dict):
        """Large bullet points with focus areas"""
        title = data.get('title', '')
        bullets = data.get('bullets', [])
        
        try:
            # Header
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.8,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                font_family=self.typography.font_family_heading
            )
            
            # Large bullet cards in 2x2 grid
            positions = [(0.6, 1.5), (6.8, 1.5), (0.6, 4.3), (6.8, 4.3)]
            
            for i, bullet in enumerate(bullets[:4]):
                left, top = positions[i]
                num = i + 1
                
                # Number badge
                badge = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(left), Inches(top),
                    Inches(0.5), Inches(0.5)
                )
                badge.fill.solid()
                badge.fill.fore_color.rgb = hex_to_rgb(self.theme.primary)
                badge.line.fill.background()
                
                self._add_text_box(
                    slide, str(num),
                    left, top + 0.08, 0.5, 0.4,
                    font_size=self.typography.body_large,
                    font_color="#FFFFFF",
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
                # Bullet text card
                self._add_glass_card(
                    slide, left + 0.7, top, 5.5, 2.5,
                    title=bullet.get('title', ''),
                    content=bullet.get('desc', '')
                )
                
        except Exception as e:
            logger.error(f"Error in bullet focus slide: {e}")
            raise
    
    def _render_metrics_slide(self, slide, data: Dict):
        """Key metrics dashboard layout"""
        title = data.get('title', '')
        metrics = data.get('metrics', [])
        
        try:
            # Header
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.8,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                font_family=self.typography.font_family_heading
            )
            
            if not metrics:
                return
            
            # Metrics grid
            num_metrics = len(metrics)
            
            # Calculate grid layout
            if num_metrics <= 3:
                cols = num_metrics
                rows = 1
            else:
                cols = 3
                rows = 2
            
            card_width = 3.8
            card_height = 2.5 if rows == 2 else 5.0
            gap = 0.4
            start_top = 1.5
            
            for i, metric in enumerate(metrics[:6]):
                row = i // cols
                col = i % cols
                
                # Center the grid
                total_width = (cols * card_width) + ((cols - 1) * gap)
                start_left = (13.333 - total_width) / 2
                
                left = start_left + (col * (card_width + gap))
                top = start_top + (row * (card_height + gap))
                
                # Metric card with large number
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(left), Inches(top),
                    Inches(card_width), Inches(card_height)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = hex_to_rgb(self.theme.surface)
                card.line.color.rgb = hex_to_rgb(self.theme.border)
                
                # Large metric value
                value = metric.get('value', '0')
                self._add_text_box(
                    slide, value,
                    left, top + 0.3, card_width, 1.0,
                    font_size=self.typography.hero if rows == 1 else self.typography.h1,
                    font_color=self.theme.primary,
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
                # Label
                label = metric.get('label', '')
                self._add_text_box(
                    slide, label,
                    left, top + 1.2 if rows == 1 else top + 1.5, 
                    card_width, 0.5,
                    font_size=self.typography.body,
                    font_color=self.theme.text_secondary,
                    align=PP_ALIGN.CENTER
                )
                
                # Trend indicator if provided
                trend = metric.get('trend')
                if trend:
                    trend_color = self.theme.success if trend.startswith('+') else self.theme.error
                    self._add_text_box(
                        slide, trend,
                        left, top + 1.8 if rows == 1 else top + 2.0, 
                        card_width, 0.4,
                        font_size=self.typography.caption,
                        font_color=trend_color,
                        bold=True,
                        align=PP_ALIGN.CENTER
                    )
                    
        except Exception as e:
            logger.error(f"Error in metrics slide: {e}")
            raise
    
    def _render_grid_slide(self, slide, data: Dict):
        """Flexible grid layout for multiple items"""
        title = data.get('title', '')
        items = data.get('items', [])
        columns = data.get('columns', 3)
        
        try:
            # Header
            if title:
                self._add_text_box(
                    slide, title,
                    0.6, 0.5, 12.1, 0.8,
                    font_size=self.typography.h1,
                    font_color=self.theme.text_primary,
                    bold=True,
                    font_family=self.typography.font_family_heading
                )
            
            # Calculate grid
            margin = 0.6
            gap = 0.3
            total_width = 13.333 - (margin * 2)
            item_width = (total_width - (gap * (columns - 1))) / columns
            item_height = 2.5
            start_top = 1.5 if title else 0.8
            
            for i, item in enumerate(items):
                row = i // columns
                col = i % columns
                
                left = margin + (col * (item_width + gap))
                top = start_top + (row * (item_height + gap))
                
                # Item card
                self._add_glass_card(
                    slide, left, top, item_width, item_height,
                    title=item.get('title', ''),
                    content=item.get('content', ''),
                    accent_color=self.theme.primary if i == 0 else None
                )
                
        except Exception as e:
            logger.error(f"Error in grid slide: {e}")
            raise
    
    def _render_hero_slide(self, slide, data: Dict):
        """Hero section with large title and call to action"""
        title = data.get('title', '')
        subtitle = data.get('subtitle', '')
        cta = data.get('cta', '')
        image_path = data.get('image', '')
        
        try:
            # Background image if provided
            if image_path and os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(0), Inches(0),
                        Inches(13.333), Inches(7.5)
                    )
                    
                    # Add overlay
                    overlay = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(13.333), Inches(7.5)
                    )
                    overlay.fill.solid()
                    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
                    overlay.fill.transparency = 0.6
                    overlay.line.fill.background()
                    
                    # Send to back
                    spTree = slide.shapes._spTree
                    sp = overlay._element
                    spTree.remove(sp)
                    spTree.insert(2, sp)
                    
                    text_color = "#FFFFFF"
                except Exception:
                    text_color = self.theme.text_primary
            else:
                text_color = self.theme.text_primary
            
            # Hero title
            self._add_text_box(
                slide, title,
                0.6, 2.0, 12.1, 2.0,
                font_size=self.typography.hero,
                font_color=text_color,
                bold=True,
                align=PP_ALIGN.CENTER,
                font_family=self.typography.font_family_heading
            )
            
            # Subtitle
            if subtitle:
                self._add_text_box(
                    slide, subtitle,
                    0.6, 3.5, 12.1, 1.0,
                    font_size=self.typography.h3,
                    font_color=text_color,
                    align=PP_ALIGN.CENTER
                )
            
            # Call to action
            if cta:
                cta_left = (13.333 - 2.0) / 2
                cta_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cta_left), Inches(4.5),
                    Inches(2.0), Inches(0.6)
                )
                cta_box.fill.solid()
                cta_box.fill.fore_color.rgb = hex_to_rgb(self.theme.primary)
                cta_box.line.fill.background()
                
                self._add_text_box(
                    slide, cta,
                    cta_left, 4.6, 2.0, 0.5,
                    font_size=self.typography.body,
                    font_color="#FFFFFF",
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
        except Exception as e:
            logger.error(f"Error in hero slide: {e}")
            raise
    
    def _render_team_slide(self, slide, data: Dict):
        """Team member showcase"""
        title = data.get('title', '')
        members = data.get('members', [])
        
        try:
            # Header
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.8,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                align=PP_ALIGN.CENTER,
                font_family=self.typography.font_family_heading
            )
            
            # Member cards
            margin = 0.6
            card_width = 2.5
            card_height = 4.0
            gap = (13.333 - margin * 2 - card_width * len(members[:4])) / (len(members[:4]) - 1) if len(members[:4]) > 1 else 0
            top = 1.8
            
            for i, member in enumerate(members[:4]):
                left = margin + i * (card_width + gap)
                
                # Photo placeholder
                photo = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(left + 0.5), Inches(top),
                    Inches(1.5), Inches(1.5)
                )
                photo.fill.solid()
                photo.fill.fore_color.rgb = hex_to_rgb(self.theme.surface_variant)
                photo.line.color.rgb = hex_to_rgb(self.theme.border)
                
                # Member info
                self._add_text_box(
                    slide, member.get('name', ''),
                    left, top + 1.7, card_width, 0.4,
                    font_size=self.typography.h4,
                    font_color=self.theme.text_primary,
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
                self._add_text_box(
                    slide, member.get('role', ''),
                    left, top + 2.2, card_width, 0.4,
                    font_size=self.typography.body,
                    font_color=self.theme.text_secondary,
                    align=PP_ALIGN.CENTER
                )
                
                self._add_text_box(
                    slide, member.get('bio', ''),
                    left, top + 2.7, card_width, 1.0,
                    font_size=self.typography.caption,
                    font_color=self.theme.text_muted,
                    align=PP_ALIGN.CENTER
                )
                
        except Exception as e:
            logger.error(f"Error in team slide: {e}")
            raise
    
    def _render_process_slide(self, slide, data: Dict):
        """Process flow visualization"""
        title = data.get('title', '')
        steps = data.get('steps', [])
        
        try:
            # Header
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.8,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                font_family=self.typography.font_family_heading
            )
            
            if not steps:
                return
            
            # Process diagram
            margin = 0.6
            step_width = (13.333 - margin * 2) / len(steps[:5])
            top = 2.0
            height = 3.5
            
            for i, step in enumerate(steps[:5]):
                left = margin + i * step_width
                
                # Step circle
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(left + step_width/2 - 0.4), Inches(top),
                    Inches(0.8), Inches(0.8)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = hex_to_rgb(self.theme.primary)
                circle.line.fill.background()
                
                # Step number
                self._add_text_box(
                    slide, str(i + 1),
                    left + step_width/2 - 0.4, top + 0.15, 0.8, 0.5,
                    font_size=self.typography.h4,
                    font_color="#FFFFFF",
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
                # Step title
                self._add_text_box(
                    slide, step.get('title', ''),
                    left, top + 1.0, step_width, 0.5,
                    font_size=self.typography.h4,
                    font_color=self.theme.text_primary,
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
                # Step description
                self._add_text_box(
                    slide, step.get('desc', ''),
                    left, top + 1.6, step_width, 1.5,
                    font_size=self.typography.caption,
                    font_color=self.theme.text_secondary,
                    align=PP_ALIGN.CENTER,
                    line_spacing=self.typography.line_height_body
                )
                
                # Connector arrow (except for last step)
                if i < len(steps[:5]) - 1:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        Inches(left + step_width - 0.3), Inches(top + 0.3),
                        Inches(0.4), Inches(0.4)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = hex_to_rgb(self.theme.text_muted)
                    arrow.line.fill.background()
                    
        except Exception as e:
            logger.error(f"Error in process slide: {e}")
            raise
    
    def _render_stats_slide(self, slide, data: Dict):
        """Statistical dashboard with visual elements"""
        title = data.get('title', '')
        stats = data.get('stats', [])
        
        try:
            # Header
            self._add_text_box(
                slide, title,
                0.6, 0.5, 12.1, 0.8,
                font_size=self.typography.h1,
                font_color=self.theme.text_primary,
                bold=True,
                font_family=self.typography.font_family_heading
            )
            
            if not stats:
                return
            
            # Stats grid
            margin = 0.6
            card_width = 3.8
            card_height = 3.0
            gap = 0.4
            start_top = 1.5
            
            num_stats = len(stats[:4])
            cols = min(num_stats, 2)
            rows = (num_stats + cols - 1) // cols
            
            total_width = (cols * card_width) + ((cols - 1) * gap)
            start_left = (13.333 - total_width) / 2
            
            for i, stat in enumerate(stats[:4]):
                row = i // cols
                col = i % cols
                
                left = start_left + (col * (card_width + gap))
                top = start_top + (row * (card_height + gap))
                
                # Stat card with visual element
                card = self._add_glass_card(
                    slide, left, top, card_width, card_height,
                    title=stat.get('label', ''),
                    content="",
                    rounded=True,
                    shadow=True
                )
                
                # Large stat value
                value = stat.get('value', '0')
                self._add_text_box(
                    slide, value,
                    left, top + 0.8, card_width, 1.0,
                    font_size=self.typography.hero,
                    font_color=self.theme.primary,
                    bold=True,
                    align=PP_ALIGN.CENTER
                )
                
                # Change indicator
                change = stat.get('change')
                if change:
                    change_color = self.theme.success if change.startswith('+') else self.theme.error
                    self._add_text_box(
                        slide, change,
                        left, top + 1.8, card_width, 0.4,
                        font_size=self.typography.body,
                        font_color=change_color,
                        bold=True,
                        align=PP_ALIGN.CENTER
                    )
                
                # Context
                context = stat.get('context')
                if context:
                    self._add_text_box(
                        slide, context,
                        left, top + 2.3, card_width, 0.4,
                        font_size=self.typography.caption,
                        font_color=self.theme.text_muted,
                        align=PP_ALIGN.CENTER
                    )
                    
        except Exception as e:
            logger.error(f"Error in stats slide: {e}")
            raise
    
    def _add_slide_number(self, slide):
        """Add subtle slide number"""
        try:
            self._add_text_box(
                slide, str(self.current_slide_num),
                12.5, 7.0, 0.5, 0.3,
                font_size=self.typography.caption,
                font_color=self.theme.text_muted,
                align=PP_ALIGN.RIGHT
            )
        except Exception as e:
            logger.debug(f"Failed to add slide number: {e}")
    
    def _add_transitions(self):
        """Add subtle fade transitions between slides"""
        # Note: python-pptx has limited transition support
        # This is a placeholder for future implementation
        pass
    
    def get_report(self) -> Dict:
        """Get generation report"""
        return {
            "total_slides": self.total_slides,
            "generated_slides": self.current_slide_num,
            "errors": self.errors,
            "warnings": self.warnings,
            "theme": self.theme.name,
            "timestamp": datetime.now().isoformat()
        }

# ============================================================================
# ADVANCED JSON SCHEMA WITH NEW LAYOUTS
# ============================================================================

EXAMPLE_JSON = {
    "theme": {
        "name": "light",
        "primary": "#3B82F6",
        "secondary": "#8B5CF6",
        "accent": "#F59E0B",
        "typography": {
            "font_family": "Calibri",
            "font_family_heading": "Calibri Light",
            "hero": 54,
            "h1": 40
        },
        "spacing_scale": 1.0
    },
    "settings": {
        "slide_numbers": True,
        "transitions": False,
        "strict_mode": False
    },
    "slides": [
        {
            "layout": "hero",
            "title": "Strategic Vision 2026",
            "subtitle": "Innovation Through Design Excellence",
            "cta": "Learn More",
            "image": "background.jpg"
        },
        {
            "layout": "title",
            "title": "Strategic Roadmap 2026",
            "subtitle": "Innovation Through Design Excellence",
            "author": "Design Team",
            "date": "March 2026"
        },
        {
            "layout": "content",
            "title": "Executive Summary",
            "content": "This presentation outlines our strategic vision for the upcoming fiscal year, focusing on three core pillars: innovation, sustainability, and growth.",
            "bullets": [
                "Data-driven approach ensures measurable outcomes",
                "Sustainable competitive advantages",
                "Cross-functional collaboration"
            ]
        },
        {
            "layout": "three-column-headed",
            "title": "Strategic Pillars",
            "columns": [
                {
                    "title": "Innovation",
                    "content": "• AI Integration\n• Process Automation\n• R&D Investment\n• Digital Transformation"
                },
                {
                    "title": "Sustainability",
                    "content": "• Carbon Neutral Goals\n• Green Operations\n• Circular Economy\n• ESG Compliance"
                },
                {
                    "title": "Growth",
                    "content": "• Market Expansion\n• New Revenue Streams\n• Strategic Partnerships\n• Customer Retention"
                }
            ]
        },
        {
            "layout": "metrics",
            "title": "Key Performance Indicators",
            "metrics": [
                {"value": "125%", "label": "Revenue Growth", "trend": "+15% YoY"},
                {"value": "50K+", "label": "New Customers", "trend": "+22% YoY"},
                {"value": "98%", "label": "Satisfaction", "trend": "+3% YoY"},
                {"value": "24", "label": "Markets", "trend": "+5 New"}
            ]
        },
        {
            "layout": "team",
            "title": "Leadership Team",
            "members": [
                {"name": "Sarah Chen", "role": "CEO", "bio": "15+ years in tech leadership"},
                {"name": "Michael Rodriguez", "role": "CTO", "bio": "AI and cloud expert"},
                {"name": "Emily Watson", "role": "CDO", "bio": "Design thinking advocate"},
                {"name": "James Kim", "role": "CFO", "bio": "Strategic finance leader"}
            ]
        },
        {
            "layout": "process",
            "title": "Our Process",
            "steps": [
                {"title": "Discover", "desc": "Research and analysis"},
                {"title": "Design", "desc": "Concept development"},
                {"title": "Develop", "desc": "Implementation"},
                {"title": "Deploy", "desc": "Launch and iterate"},
                {"title": "Scale", "desc": "Growth and optimization"}
            ]
        },
        {
            "layout": "stats",
            "title": "Company Stats",
            "stats": [
                {"label": "Revenue", "value": "$50M", "change": "+25%", "context": "FY 2025"},
                {"label": "Customers", "value": "10K+", "change": "+40%", "context": "Global"},
                {"label": "Markets", "value": "15", "change": "+5", "context": "Countries"},
                {"label": "Team", "value": "200+", "change": "+30%", "context": "Employees"}
            ]
        },
        {
            "layout": "image-right",
            "title": "Technology Stack",
            "content": "Our modern infrastructure leverages cloud-native architectures, microservices, and advanced analytics to deliver scalable solutions with 99.9% uptime.",
            "image": "tech_diagram.png"
        },
        {
            "layout": "timeline",
            "title": "Implementation Roadmap",
            "events": [
                {"date": "Q1 2026", "title": "Planning", "desc": "Strategy finalization"},
                {"date": "Q2 2026", "title": "Development", "desc": "Core platform build"},
                {"date": "Q3 2026", "title": "Launch", "desc": "Market release"},
                {"date": "Q4 2026", "title": "Scale", "desc": "Global expansion"}
            ]
        },
        {
            "layout": "comparison",
            "title": "Competitive Analysis",
            "left_title": "Our Solution",
            "right_title": "Traditional",
            "left_points": ["Real-time analytics", "Cloud-native", "API-first", "Scalable"],
            "right_points": ["Batch processing", "On-premise", "Monolithic", "Limited scale"]
        },
        {
            "layout": "quote",
            "quote": "Design is not just what it looks like and feels like. Design is how it works.",
            "author": "Steve Jobs",
            "role": "Co-founder, Apple Inc."
        },
        {
            "layout": "grid",
            "title": "Key Features",
            "columns": 3,
            "items": [
                {"title": "Feature 1", "content": "Description here"},
                {"title": "Feature 2", "content": "Description here"},
                {"title": "Feature 3", "content": "Description here"},
                {"title": "Feature 4", "content": "Description here"},
                {"title": "Feature 5", "content": "Description here"}
            ]
        }
    ]
}

# ============================================================================
# COMMAND LINE INTERFACE
# ============================================================================

def main():
    """Command line interface for presentation generator"""
    import argparse
    
    parser = argparse.ArgumentParser(description="Modern PowerPoint Presentation Generator")
    parser.add_argument("-i", "--input", help="Input JSON file path")
    parser.add_argument("-o", "--output", help="Output PPTX file path")
    parser.add_argument("-t", "--theme", choices=["light", "dark", "corporate", "minimalist"], 
                       help="Theme to use")
    parser.add_argument("--verbose", action="store_true", help="Enable verbose output")
    
    args = parser.parse_args()
    
    # Set log level
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    print("=" * 60)
    print("MODERN POWERPOINT GENERATOR")
    print("=" * 60)
    
    try:
        if args.input:
            # Load from file
            with open(args.input, 'r', encoding='utf-8') as f:
                json_data = json.load(f)
            
            # Override theme if specified
            if args.theme:
                json_data['theme'] = json_data.get('theme', {})
                json_data['theme']['name'] = args.theme
            
            generator = ModernPresentationGenerator(json_data)
            output_path = generator.generate(args.output)
            print(f"\n✓ Presentation generated: {output_path}")
            
            # Show report
            report = generator.get_report()
            if report['warnings']:
                print(f"\n⚠ Warnings ({len(report['warnings'])}):")
                for warning in report['warnings'][:5]:
                    print(f"  • {warning}")
            if report['errors']:
                print(f"\n✗ Errors ({len(report['errors'])}):")
                for error in report['errors'][:5]:
                    print(f"  • {error}")
                    
        else:
            # Use example
            generator = ModernPresentationGenerator(EXAMPLE_JSON)
            generator.generate("Modern_Presentation.pptx")
            
            print("\n" + "=" * 60)
            print("Available Layout Types:")
            for layout in LayoutType:
                print(f"  • {layout.value}")
            print("=" * 60)
            
    except FileNotFoundError:
        print(f"Error: File '{args.input}' not found")
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON format - {e}")
        sys.exit(1)
    except PresentationError as e:
        print(f"Error: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"Unexpected error: {e}")
        if args.verbose:
            import traceback
            traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    import sys
    main()